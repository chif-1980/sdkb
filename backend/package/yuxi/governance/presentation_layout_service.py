from __future__ import annotations

import io
import re
from collections.abc import Iterable, Sequence
from dataclasses import dataclass
from difflib import SequenceMatcher
from typing import Any

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from yuxi.services.file_preview import convert_office_to_pdf


def _clean_text(value: str | None) -> str:
    lines = [re.sub(r"\s+", " ", line).strip() for line in (value or "").splitlines()]
    return "\n".join(line for line in lines if line).strip()


def _normalized_text(value: str | None) -> str:
    return re.sub(r"[^0-9a-z\u3400-\u9fff]+", "", (value or "").lower())


def _shape_text(shape: Any) -> str:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return _clean_text("\n".join(_shape_text(child) for child in shape.shapes))
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            values = [_clean_text(cell.text) for cell in row.cells]
            if any(values):
                rows.append(" | ".join(values))
        return _clean_text("\n".join(rows))
    if getattr(shape, "has_text_frame", False):
        return _clean_text(shape.text)
    return ""


def _percentage(value: int | None, total: int) -> float:
    if not value or total <= 0:
        return 0.0
    return round(max(0.0, min(100.0, value / total * 100)), 4)


def _segment_value(segment: Any, name: str, default: Any = None) -> Any:
    if isinstance(segment, dict):
        return segment.get(name, default)
    return getattr(segment, name, default)


def _matching_segment_ids(fragment_text: str, segments: Sequence[Any]) -> list[str]:
    fragment = _normalized_text(fragment_text)
    if len(fragment) < 2:
        return []
    scored: list[tuple[float, int, str]] = []
    for segment in segments:
        segment_id = str(_segment_value(segment, "segment_id", ""))
        candidate = _normalized_text(str(_segment_value(segment, "content", "")))
        if not segment_id or len(candidate) < 2:
            continue
        if fragment in candidate or candidate in fragment:
            score = min(len(fragment), len(candidate)) / max(len(fragment), len(candidate))
            if score >= 0.25:
                scored.append((score + 1.0, -int(_segment_value(segment, "segment_index", 0)), segment_id))
            continue
        score = SequenceMatcher(None, fragment[:2000], candidate[:2000], autojunk=False).ratio()
        if score >= 0.42:
            scored.append((score, -int(_segment_value(segment, "segment_index", 0)), segment_id))
    scored.sort(reverse=True)
    return [segment_id for _score, _index, segment_id in scored[:3]]


@dataclass(frozen=True, slots=True)
class PresentationLayout:
    slide_width: int
    slide_height: int
    slides: list[dict[str, Any]]

    def as_dict(self) -> dict[str, Any]:
        return {
            "supported": True,
            "slide_count": len(self.slides),
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "aspect_ratio": round(self.slide_width / self.slide_height, 6) if self.slide_height else None,
            "slides": self.slides,
        }


def extract_pptx_layout(content: bytes, *, segments: Sequence[Any] = ()) -> PresentationLayout:
    presentation = Presentation(io.BytesIO(content))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    slides: list[dict[str, Any]] = []
    fragment_number = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        fragments = []
        for shape_number, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            if not text:
                continue
            fragment_number += 1
            left = _percentage(getattr(shape, "left", 0), slide_width)
            top = _percentage(getattr(shape, "top", 0), slide_height)
            width = _percentage(getattr(shape, "width", 0), slide_width)
            height = _percentage(getattr(shape, "height", 0), slide_height)
            width = min(width, round(100.0 - left, 4))
            height = min(height, round(100.0 - top, 4))
            fragments.append(
                {
                    "fragment_id": f"slide-{slide_number}-shape-{shape_number}",
                    "fragment_number": fragment_number,
                    "shape_number": shape_number,
                    "content": text,
                    "left": left,
                    "top": top,
                    "width": max(width, 0.5),
                    "height": max(height, 0.5),
                    "source_segment_ids": _matching_segment_ids(text, segments),
                }
            )
        slides.append(
            {
                "slide_number": slide_number,
                "fragment_count": len(fragments),
                "fragments": fragments,
            }
        )
    return PresentationLayout(slide_width=slide_width, slide_height=slide_height, slides=slides)


def render_pdf_slide(pdf_content: bytes, *, slide_number: int) -> bytes:
    with fitz.open(stream=pdf_content, filetype="pdf") as document:
        if slide_number < 1 or slide_number > document.page_count:
            raise IndexError("Slide number is out of range")
        page = document.load_page(slide_number - 1)
        target_width = 1600
        scale = min(3.0, max(1.0, target_width / max(page.rect.width, 1)))
        pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        return pixmap.tobytes("jpeg", jpg_quality=88)


async def render_pptx_slide(
    content: bytes,
    *,
    filename: str,
    slide_number: int,
    pdf_content: bytes | None = None,
) -> bytes:
    """Render one slide, optionally reusing a PDF produced for the same deck.

    Converting a PPTX with LibreOffice is much more expensive than rendering a
    page from the resulting PDF. Callers that render multiple slides should
    convert once and pass ``pdf_content`` for subsequent pages.
    """
    if pdf_content is None:
        pdf_content = await convert_office_to_pdf(filename, content)
    return render_pdf_slide(pdf_content, slide_number=slide_number)


def iter_fragment_text(layout: PresentationLayout) -> Iterable[str]:
    for slide in layout.slides:
        for fragment in slide["fragments"]:
            yield fragment["content"]
