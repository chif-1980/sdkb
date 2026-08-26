import io

import fitz
import pytest
from pptx import Presentation
from pptx.util import Inches

from yuxi.governance import presentation_layout_service
from yuxi.governance.presentation_layout_service import extract_pptx_layout, render_pptx_slide


def _presentation_bytes() -> bytes:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = first.shapes.add_textbox(Inches(1), Inches(0.6), Inches(5), Inches(0.8))
    title.text = "公司简介"
    body = first.shapes.add_textbox(Inches(1), Inches(1.8), Inches(7), Inches(2))
    body.text = "公司为客户提供咨询、实施、交付和持续运营服务。"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_title = second.shapes.add_textbox(Inches(1), Inches(0.6), Inches(5), Inches(0.8))
    second_title.text = "产品能力"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_extract_pptx_layout_keeps_real_pages_geometry_and_segment_mapping():
    layout = extract_pptx_layout(
        _presentation_bytes(),
        segments=[
            {
                "segment_id": "segment-company",
                "segment_index": 3,
                "content": "公司为客户提供咨询、实施、交付和持续运营服务。",
            }
        ],
    )

    assert len(layout.slides) == 2
    assert layout.slides[0]["slide_number"] == 1
    assert [item["content"] for item in layout.slides[0]["fragments"]] == [
        "公司简介",
        "公司为客户提供咨询、实施、交付和持续运营服务。",
    ]
    body = layout.slides[0]["fragments"][1]
    assert 0 < body["left"] < 100
    assert 0 < body["top"] < 100
    assert body["source_segment_ids"] == ["segment-company"]
    assert layout.as_dict()["slide_count"] == 2


@pytest.mark.asyncio
async def test_render_pptx_slide_returns_png(monkeypatch):
    document = fitz.open()
    page = document.new_page(width=960, height=540)
    page.insert_text((72, 72), "Slide one")
    pdf_content = document.tobytes()
    document.close()

    async def fake_convert(_filename, _content):
        return pdf_content

    monkeypatch.setattr(presentation_layout_service, "convert_office_to_pdf", fake_convert)

    image = await render_pptx_slide(b"pptx", filename="intro.pptx", slide_number=1)

    assert image.startswith(b"\x89PNG\r\n\x1a\n")


@pytest.mark.asyncio
async def test_render_pptx_slide_rejects_unknown_page(monkeypatch):
    document = fitz.open()
    document.new_page()
    pdf_content = document.tobytes()
    document.close()

    async def fake_convert(_filename, _content):
        return pdf_content

    monkeypatch.setattr(presentation_layout_service, "convert_office_to_pdf", fake_convert)

    with pytest.raises(IndexError, match="out of range"):
        await render_pptx_slide(b"pptx", filename="intro.pptx", slide_number=2)


@pytest.mark.asyncio
async def test_render_pptx_slide_reuses_preconverted_pdf(monkeypatch):
    document = fitz.open()
    document.new_page()
    pdf_content = document.tobytes()
    document.close()

    async def fail_convert(_filename, _content):
        raise AssertionError("should not convert when a cached PDF is provided")

    monkeypatch.setattr(presentation_layout_service, "convert_office_to_pdf", fail_convert)

    image = await render_pptx_slide(
        b"pptx",
        filename="intro.pptx",
        slide_number=1,
        pdf_content=pdf_content,
    )

    assert image.startswith(b"\x89PNG\r\n\x1a\n")
